import os
from pptx import Presentation

def create_ppt_from_images(image_dir, output_filename):
  prs = Presentation()
  slide_layout = prs.slide_layouts[6]  

  for filename in os.listdir(image_dir):
    filepath = os.path.join(image_dir, filename)
    if os.path.isfile(filepath) and filename.lower().endswith(('.png', '.jpg', '.jpeg')):
      slide = prs.slides.add_slide(slide_layout)
      pic = slide.shapes.add_picture(filepath, 0, 0, width=prs.slide_width, height=prs.slide_height)

  prs.save(output_filename)

if __name__ == "__main__":
  image_dir = "Path to folder
  output_filename = "images.pptx"

  create_ppt_from_images(image_dir, output_filename)
  print(f"PowerPoint presentation created: {output_filename}")